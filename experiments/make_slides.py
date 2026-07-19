from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# Color palette
ORANGE      = RGBColor(0xE5, 0x5A, 0x00)   # deep orange
ORANGE_LIGHT= RGBColor(0xFF, 0x8C, 0x2F)   # lighter orange accent
DARK        = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF4, 0xF4, 0xF4)
MID_GRAY    = RGBColor(0x88, 0x88, 0x88)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=6, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_bullet_box(slide, items, l, t, w, h,
                   font_size=18, color=WHITE, dot_color=None,
                   line_height=6, bold_first=False):
    """Items is list of strings. Adds a textbox with bullet-like layout."""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(line_height)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold_first and (item == items[0])
    return txBox


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)

# Full dark background
add_rect(s1, 0, 0, W, H, fill_rgb=DARK)

# Orange left accent bar
add_rect(s1, 0, 0, Inches(0.45), H, fill_rgb=ORANGE)

# Orange bottom stripe
add_rect(s1, 0, H - Inches(0.9), W, Inches(0.9), fill_rgb=ORANGE)

# Decorative top-right corner box (lighter orange)
add_rect(s1, W - Inches(3.5), 0, Inches(3.5), Inches(1.2), fill_rgb=ORANGE_LIGHT)

# Main title
add_text(s1,
         "AI within Course Topics:",
         Inches(0.8), Inches(1.8), Inches(10), Inches(1.2),
         font_size=44, bold=True, color=WHITE)
add_text(s1,
         "Mapping the ML Curriculum",
         Inches(0.8), Inches(2.85), Inches(10), Inches(1.1),
         font_size=44, bold=True, color=ORANGE_LIGHT)

# Divider line
add_rect(s1, Inches(0.8), Inches(4.1), Inches(5), Inches(0.05), fill_rgb=ORANGE_LIGHT)

# Authors
add_text(s1,
         "Alina G.  ·  Shrivarshini N.",
         Inches(0.8), Inches(4.3), Inches(8), Inches(0.6),
         font_size=20, color=LIGHT_GRAY)

# Subtitle / date
add_text(s1,
         "CIC Retreat  ·  June 2026",
         Inches(0.8), Inches(4.95), Inches(8), Inches(0.5),
         font_size=16, color=MID_GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — BIG QUESTIONS  (Shrivarshini's slides)
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, W, H, fill_rgb=DARK)
add_rect(s2, 0, 0, W, Inches(1.5), fill_rgb=ORANGE)          # top header bar
add_rect(s2, 0, H - Inches(0.35), W, Inches(0.35), fill_rgb=ORANGE)  # bottom bar

# Slide header
add_text(s2, "The Big Questions", Inches(0.55), Inches(0.25), Inches(9), Inches(0.9),
         font_size=36, bold=True, color=WHITE)

# Motivating line
add_text(s2,
         "As AI reshapes every field, what exactly are students learning — and is it enough?",
         Inches(0.55), Inches(1.65), Inches(12.2), Inches(0.65),
         font_size=16, italic=True, color=LIGHT_GRAY)

# 4 question cards
questions = [
    ("01", "What's being taught?",
     "Which ML topics appear across university courses, and at what level of granularity?"),
    ("02", "How much time?",
     "How many weeks does each topic receive, and how does emphasis vary by school?"),
    ("03", "Are there discrepancies?",
     "Do courses differ significantly in scope, depth, or ordering of material?"),
    ("04", "Are there gaps?",
     "What's missing from current curricula — and how much does it matter?"),
]

card_w = Inches(2.85)
card_h = Inches(3.7)
card_top = Inches(2.45)
gap = Inches(0.28)

for i, (num, q, detail) in enumerate(questions):
    cl = Inches(0.45) + i * (card_w + gap)
    # card background
    add_rect(s2, cl, card_top, card_w, card_h,
             fill_rgb=RGBColor(0x28, 0x28, 0x42), line_rgb=ORANGE_LIGHT, line_pt=1.2)
    # number badge
    add_rect(s2, cl + Inches(0.18), card_top + Inches(0.18),
             Inches(0.55), Inches(0.42), fill_rgb=ORANGE)
    add_text(s2, num,
             cl + Inches(0.18), card_top + Inches(0.16),
             Inches(0.55), Inches(0.44),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # question heading
    add_text(s2, q,
             cl + Inches(0.18), card_top + Inches(0.75),
             card_w - Inches(0.36), Inches(0.85),
             font_size=16, bold=True, color=ORANGE_LIGHT)
    # detail text
    add_text(s2, detail,
             cl + Inches(0.18), card_top + Inches(1.65),
             card_w - Inches(0.36), Inches(1.75),
             font_size=13, color=LIGHT_GRAY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE APPROACH
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, W, H, fill_rgb=LIGHT_GRAY)
add_rect(s3, 0, 0, W, Inches(1.5), fill_rgb=ORANGE)
add_rect(s3, 0, H - Inches(0.35), W, Inches(0.35), fill_rgb=DARK)

add_text(s3, "The Approach", Inches(0.55), Inches(0.25), Inches(9), Inches(0.9),
         font_size=36, bold=True, color=WHITE)

steps = [
    ("1", "Standardize", "Build taxonomy",
     "Adopted the ACM 2023 CS Curricula taxonomy — two levels:\nL1 (broad categories) and L2 (fine-grained subtopics)"),
    ("2", "Collect", "Gather syllabuses",
     "Identified 500+ ML courses at US & Canadian universities;\nselected those with week-by-week schedules"),
    ("3", "Label", "Extract topics",
     "Manually labeled 50 syllabuses; built an LLM-based pipeline\nto automate extraction and compare against human labels"),
]

box_w = Inches(3.6)
box_h = Inches(4.2)
box_top = Inches(1.85)
box_gap = Inches(0.45)

for i, (num, title, sub, detail) in enumerate(steps):
    bl = Inches(0.5) + i * (box_w + box_gap)
    add_rect(s3, bl, box_top, box_w, box_h,
             fill_rgb=WHITE, line_rgb=ORANGE, line_pt=1.5)
    # colored number circle area
    add_rect(s3, bl, box_top, box_w, Inches(0.95), fill_rgb=ORANGE)
    add_text(s3, f"Step {num}", bl + Inches(0.15), box_top + Inches(0.1),
             box_w - Inches(0.3), Inches(0.7),
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s3, title,
             bl + Inches(0.15), box_top + Inches(1.1),
             box_w - Inches(0.3), Inches(0.65),
             font_size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s3, sub,
             bl + Inches(0.15), box_top + Inches(1.7),
             box_w - Inches(0.3), Inches(0.5),
             font_size=14, bold=False, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)
    add_text(s3, detail,
             bl + Inches(0.2), box_top + Inches(2.3),
             box_w - Inches(0.4), Inches(1.7),
             font_size=13, color=DARK, wrap=True)

# Arrows between boxes
for i in range(2):
    ax = Inches(0.5) + (i + 1) * (box_w + box_gap) - box_gap + Inches(0.05)
    ay = box_top + box_h / 2 - Inches(0.15)
    add_text(s3, "▶", ax - Inches(0.28), ay, Inches(0.4), Inches(0.4),
             font_size=20, color=ORANGE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT WE'VE DONE SO FAR
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, W, H, fill_rgb=DARK)
add_rect(s4, 0, 0, W, Inches(1.5), fill_rgb=ORANGE)
add_rect(s4, 0, H - Inches(0.35), W, Inches(0.35), fill_rgb=ORANGE)

add_text(s4, "What We've Done So Far", Inches(0.55), Inches(0.25), Inches(10), Inches(0.9),
         font_size=36, bold=True, color=WHITE)

milestones = [
    ("Refined taxonomy",     "Aligned topic labels with ACM 2023 L1/L2 structure across the team"),
    ("500+ courses scanned", "Built a comprehensive list of potential ML courses at US & Canadian universities"),
    ("50 syllabuses labeled","Manually labeled week-by-week schedules; includes 5 Canadian schools"),
    ("Auto-labeling tool",   "Developed an LLM pipeline (DeepSeek) to extract topics automatically"),
    ("Baseline comparison",  "Evaluated against human labels using Precision@O, Recall@O, and F1@O metrics"),
]

row_h   = Inches(0.9)
row_gap = Inches(0.18)
row_top = Inches(1.7)

for i, (heading, detail) in enumerate(milestones):
    rt = row_top + i * (row_h + row_gap)
    # row background (alternating slight shade)
    bg = RGBColor(0x28, 0x28, 0x42) if i % 2 == 0 else RGBColor(0x22, 0x22, 0x38)
    add_rect(s4, Inches(0.45), rt, W - Inches(0.9), row_h, fill_rgb=bg)
    # orange left tick
    add_rect(s4, Inches(0.45), rt, Inches(0.12), row_h, fill_rgb=ORANGE)
    # check mark
    add_text(s4, "✓", Inches(0.6), rt + Inches(0.08), Inches(0.5), row_h - Inches(0.1),
             font_size=20, bold=True, color=ORANGE)
    # heading
    add_text(s4, heading, Inches(1.2), rt + Inches(0.08),
             Inches(3.5), row_h - Inches(0.1),
             font_size=16, bold=True, color=WHITE)
    # detail
    add_text(s4, detail, Inches(4.8), rt + Inches(0.1),
             Inches(8.0), row_h - Inches(0.15),
             font_size=14, color=LIGHT_GRAY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES 5-8 — RESULTS PLACEHOLDER SLIDES
# (Alina will drop in the actual charts)
# ═══════════════════════════════════════════════════════════════════════════════

chart_slides = [
    ("Average Topics: True vs. Extracted",
     "LLM Powered Approach",
     "The LLM extracts ~12 topics on average vs. ~28 in the human labels —\nsuggesting the model is conservative rather than noisy."),
    ("F1@O Score Distribution",
     "LLM Powered Approach",
     "Most syllabuses score between 0.5–0.7 F1, indicating moderate\nbut consistent overlap with human-labeled topics."),
    ("Precision@O Distribution",
     "LLM Powered Approach",
     "Bimodal pattern: when the LLM picks a topic, it's almost always correct.\nThe challenge is coverage, not accuracy."),
    ("Recall@O Distribution",
     "LLM Powered Approach",
     "Recall is spread 0–1, peaking around 0.5–0.6 — the LLM misses\nroughly half of the true topics per syllabus."),
]

for title, subtitle, insight in chart_slides:
    sc = prs.slides.add_slide(blank_layout)
    add_rect(sc, 0, 0, W, H, fill_rgb=LIGHT_GRAY)
    add_rect(sc, 0, 0, W, Inches(1.3), fill_rgb=DARK)
    add_rect(sc, 0, H - Inches(0.35), W, Inches(0.35), fill_rgb=ORANGE)

    add_text(sc, title, Inches(0.55), Inches(0.18), Inches(10), Inches(0.75),
             font_size=28, bold=True, color=WHITE)
    add_text(sc, subtitle, Inches(0.55), Inches(0.88), Inches(10), Inches(0.35),
             font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC), italic=True)

    # Chart placeholder box
    add_rect(sc, Inches(0.6), Inches(1.5), Inches(8.8), Inches(4.5),
             fill_rgb=WHITE, line_rgb=ORANGE, line_pt=1.2)
    add_text(sc, "[ chart goes here ]",
             Inches(0.6), Inches(3.2), Inches(8.8), Inches(1.0),
             font_size=18, color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Insight callout box (right side)
    add_rect(sc, Inches(9.7), Inches(1.5), Inches(3.3), Inches(4.5),
             fill_rgb=DARK, line_rgb=ORANGE_LIGHT, line_pt=1.5)
    add_text(sc, "Key Takeaway",
             Inches(9.85), Inches(1.65), Inches(3.0), Inches(0.55),
             font_size=14, bold=True, color=ORANGE_LIGHT)
    add_text(sc, insight,
             Inches(9.85), Inches(2.3), Inches(3.0), Inches(3.3),
             font_size=12, color=LIGHT_GRAY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE — NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════════════
sn = prs.slides.add_slide(blank_layout)
add_rect(sn, 0, 0, W, H, fill_rgb=DARK)
add_rect(sn, 0, 0, W, Inches(1.5), fill_rgb=ORANGE)
add_rect(sn, 0, H - Inches(0.35), W, Inches(0.35), fill_rgb=ORANGE)

add_text(sn, "What's Next", Inches(0.55), Inches(0.25), Inches(9), Inches(0.9),
         font_size=36, bold=True, color=WHITE)

nexts = [
    "Expand to 60+ labeled syllabuses and add university-level metadata",
    "Improve LLM recall — explore prompting strategies and chain-of-thought extraction",
    "Cross-university analysis: which topics cluster together? Which are under-taught?",
    "Integrate findings into the AAAI 2026 paper draft",
]

for i, txt in enumerate(nexts):
    nt = Inches(1.85) + i * Inches(1.1)
    add_rect(sn, Inches(0.45), nt, Inches(0.45), Inches(0.75), fill_rgb=ORANGE)
    add_text(sn, str(i + 1), Inches(0.45), nt + Inches(0.05), Inches(0.45), Inches(0.65),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sn, Inches(0.95), nt, Inches(11.8), Inches(0.75),
             fill_rgb=RGBColor(0x28, 0x28, 0x42))
    add_text(sn, txt, Inches(1.1), nt + Inches(0.1), Inches(11.5), Inches(0.6),
             font_size=16, color=LIGHT_GRAY)


# ── save ──────────────────────────────────────────────────────────────────────
out = "/Users/shrivarshininarayanan/meddata-guardian-1/experiments/ml_curriculum_slides.pptx"
prs.save(out)
print(f"Saved to {out}")
